from pptx import Presentation
import io

def extract_text_from_pptx(raw_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(raw_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)
