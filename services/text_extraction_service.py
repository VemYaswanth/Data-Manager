# services/text_extraction_service.py

import os
from io import BytesIO
from core.logger import logger


# ---------------------------
# Safe UTF-8 decode fallback
# ---------------------------
def _safe_decode(raw: bytes) -> str:
    try:
        return raw.decode("utf-8", errors="ignore")
    except Exception:
        return ""


# ---------------------------
# EXTRACTORS
# ---------------------------
def extract_pdf(raw_bytes: bytes) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=raw_bytes, filetype="pdf")
        return "\n".join([page.get_text() for page in doc])
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return ""


def extract_docx(raw_bytes: bytes) -> str:
    try:
        from docx import Document
        doc = Document(BytesIO(raw_bytes))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return ""


def extract_pptx(raw_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(BytesIO(raw_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return ""


def extract_xlsx(raw_bytes: bytes) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(BytesIO(raw_bytes), data_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(v) for v in row if v is not None]
                if row_vals:
                    texts.append(" ".join(row_vals))
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        return ""


def extract_image(raw_bytes: bytes) -> str:
    try:
        import easyocr
        import numpy as np
        import cv2

        reader = easyocr.Reader(["en"])
        np_img = np.frombuffer(raw_bytes, np.uint8)
        img = cv2.imdecode(np_img, cv2.IMREAD_COLOR)
        result = reader.readtext(img, detail=0)
        return "\n".join(result)
    except Exception as e:
        logger.error(f"OCR failed: {e}")
        return ""


def extract_html(raw_bytes: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
        html = raw_bytes.decode("utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        return soup.get_text(separator="\n")
    except Exception as e:
        logger.error(f"HTML extraction failed: {e}")
        return ""


# ---------------------------
# Main extraction router
# ---------------------------
def extract_text_from_bytes(filename: str, raw_bytes: bytes) -> str:
    ext = os.path.splitext(filename.lower())[1]

    try:
        if ext in [".txt", ".md", ".log", ".json"]:
            return _safe_decode(raw_bytes)

        if ext == ".csv":
            return _safe_decode(raw_bytes)

        if ext == ".pdf":
            return extract_pdf(raw_bytes)

        if ext == ".docx":
            return extract_docx(raw_bytes)

        if ext == ".pptx":
            return extract_pptx(raw_bytes)

        if ext in [".xlsx", ".xlsm", ".xltx"]:
            return extract_xlsx(raw_bytes)

        if ext in [".jpg", ".jpeg", ".png"]:
            return extract_image(raw_bytes)

        if ext in [".html", ".htm"]:
            return extract_html(raw_bytes)

        # Fallback attempt
        return _safe_decode(raw_bytes)

    except Exception as e:
        logger.error(f"Extraction router failed for {filename}: {e}")
        return ""
